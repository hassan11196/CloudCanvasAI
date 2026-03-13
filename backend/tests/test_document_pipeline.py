"""
Test the document generation pipeline end-to-end.

Tests:
1. Sandbox creation
2. DOCX file creation in sandbox
3. DOCX file accessibility (can be read)
4. DOCX file renderability (valid format, can be parsed)
"""

import json
import os
import socket
import sys
import tempfile
import time
import uuid
from io import BytesIO
from pathlib import Path
from urllib.parse import urlparse

import pytest

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))

from sandbox_manager import SandboxManager, SANDBOX_ROOT, SANDBOX_WORKSPACE
from session_files import session_storage_path, session_workspace_dir


def _e2b_host() -> str:
    """Resolve the E2B API host from env (supports overrides)."""
    base_url = os.getenv("E2B_API_URL")
    if not base_url:
        domain = os.getenv("E2B_DOMAIN", "e2b.app")
        base_url = f"https://api.{domain}"

    parsed = urlparse(base_url)
    return parsed.hostname or base_url


def _can_reach_e2b_api() -> bool:
    """Detect whether the E2B API is reachable (DNS + TCP)."""
    host = _e2b_host()
    if not host:
        return False

    try:
        # Quick DNS resolution
        socket.getaddrinfo(host, 443)
        # Quick TCP connectivity check
        with socket.create_connection((host, 443), timeout=3):
            return True
    except Exception:
        return False


# Skip when E2B key missing or the API host is unreachable
pytestmark = [
    pytest.mark.skipif(
        not os.getenv("E2B_API_KEY"),
        reason="E2B_API_KEY environment variable not set",
    ),
    pytest.mark.skipif(
        not _can_reach_e2b_api(),
        reason="E2B API unreachable (DNS or TCP 443 blocked)",
    ),
]


def _ensure_html2pptx_script(sandbox) -> str:
    """Ensure the html2pptx script exists in the sandbox."""
    remote_path = f"{SANDBOX_ROOT}/scripts/html2pptx.js"
    candidate_paths = [
        # Preferred path (requested skill location)
        Path(__file__).resolve().parents[2] / "skills" / "skills" / "pptx" / "scripts" / "html2pptx.js",
        # Fallback copy under backend
        Path(__file__).parent.parent / "skills" / "pptx" / "scripts" / "html2pptx.js",
    ]

    local_path = next((p for p in candidate_paths if p.exists()), None)
    if local_path is None:
        raise FileNotFoundError("Local html2pptx.js not found in expected skill paths")

    try:
        sandbox.files.make_dir(f"{SANDBOX_ROOT}/scripts")
    except Exception:
        pass  # Directory may already exist

    sandbox.files.write(remote_path, local_path.read_text())

    try:
        sandbox.files.read(remote_path)
    except Exception as e:
        raise RuntimeError(f"html2pptx.js not available in sandbox: {e}")

    return remote_path


class TestSandboxCreation:
    """Test 1: Verify sandbox can be created."""

    def test_sandbox_manager_initialization(self):
        """SandboxManager initializes without errors."""
        manager = SandboxManager()
        assert manager is not None
        assert manager.api_key is not None, "E2B_API_KEY should be set"

    def test_sandbox_creation(self):
        """New sandbox can be created for a user."""
        manager = SandboxManager()
        test_user_id = f"test-user-{uuid.uuid4()}"

        try:
            sandbox = manager.create_sandbox(test_user_id)
            assert sandbox is not None
            assert sandbox.is_running(), "Sandbox should be running after creation"

            # Verify sandbox is tracked
            retrieved = manager.get_sandbox(test_user_id)
            assert retrieved is not None
            assert retrieved.sandbox_id == sandbox.sandbox_id

        finally:
            # Cleanup
            manager.close_sandbox(test_user_id)

    def test_sandbox_workspace_exists(self):
        """Sandbox has required workspace directories."""
        manager = SandboxManager()
        test_user_id = f"test-user-{uuid.uuid4()}"

        try:
            sandbox = manager.create_sandbox(test_user_id)

            # Check workspace directory exists
            entries = sandbox.files.list(SANDBOX_WORKSPACE)
            assert entries is not None, "Workspace directory should exist"

            # Check scripts directory exists
            scripts_entries = sandbox.files.list(f"{SANDBOX_ROOT}/scripts")
            assert scripts_entries is not None, "Scripts directory should exist"

            # Verify create_docx.py exists
            script_names = [e.name for e in scripts_entries]
            assert "create_docx.py" in script_names, "create_docx.py should be in scripts"

        finally:
            manager.close_sandbox(test_user_id)


class TestDocxCreation:
    """Test 2: Verify DOCX files can be created in sandbox."""

    @pytest.fixture
    def sandbox_setup(self):
        """Create a sandbox for testing and clean up after."""
        manager = SandboxManager()
        test_user_id = f"test-user-{uuid.uuid4()}"
        session_id = str(uuid.uuid4())

        sandbox = manager.create_sandbox(test_user_id)

        # Create session workspace
        session_dir = session_workspace_dir(session_id)
        try:
            sandbox.files.make_dir(session_dir)
        except Exception:
            pass  # May already exist

        yield {
            "manager": manager,
            "sandbox": sandbox,
            "user_id": test_user_id,
            "session_id": session_id,
            "session_dir": session_dir,
        }

        # Cleanup
        manager.close_sandbox(test_user_id)

    def test_create_simple_docx(self, sandbox_setup):
        """Simple DOCX can be created using create_docx.py script."""
        sandbox = sandbox_setup["sandbox"]
        session_dir = sandbox_setup["session_dir"]

        # Document content
        doc_content = {
            "title": "Test Document",
            "sections": [
                {
                    "heading": "Introduction",
                    "level": 1,
                    "content": "This is a test document."
                }
            ]
        }

        output_path = f"{session_dir}/test_simple.docx"
        json_content = json.dumps(doc_content)

        # Run create_docx.py script
        result = sandbox.commands.run(
            f'python {SANDBOX_ROOT}/scripts/create_docx.py "{output_path}" \'{json_content}\'',
            timeout=60,
            cwd=session_dir,
        )

        assert result.exit_code == 0, f"Script failed: {result.stderr}"
        assert "Document created" in result.stdout, f"Expected success message, got: {result.stdout}"

    def test_create_complex_docx(self, sandbox_setup):
        """Complex DOCX with multiple sections can be created."""
        sandbox = sandbox_setup["sandbox"]
        session_dir = sandbox_setup["session_dir"]

        doc_content = {
            "title": "Complex Test Document",
            "sections": [
                {
                    "heading": "Introduction",
                    "level": 1,
                    "content": "This document tests all features."
                },
                {
                    "heading": "Bullet Points",
                    "level": 2,
                    "bullet_points": ["Item 1", "Item 2", "Item 3"]
                },
                {
                    "heading": "Numbered List",
                    "level": 2,
                    "numbered_list": ["First", "Second", "Third"]
                },
                {
                    "heading": "Table",
                    "level": 2,
                    "table": {
                        "headers": ["Name", "Value"],
                        "rows": [["Test", "123"], ["Sample", "456"]]
                    }
                }
            ]
        }

        output_path = f"{session_dir}/test_complex.docx"

        # Write JSON to a temp file to avoid shell escaping issues
        json_file = f"{session_dir}/content.json"
        sandbox.files.write(json_file, json.dumps(doc_content))

        result = sandbox.commands.run(
            f'python {SANDBOX_ROOT}/scripts/create_docx.py "{output_path}" --file "{json_file}"',
            timeout=60,
            cwd=session_dir,
        )

        assert result.exit_code == 0, f"Script failed: {result.stderr}"


class TestDocxAccessibility:
    """Test 3: Verify DOCX files can be read from sandbox."""

    @pytest.fixture
    def created_docx(self):
        """Create a DOCX file and return sandbox context."""
        manager = SandboxManager()
        test_user_id = f"test-user-{uuid.uuid4()}"
        session_id = str(uuid.uuid4())

        sandbox = manager.create_sandbox(test_user_id)
        session_dir = session_workspace_dir(session_id)

        try:
            sandbox.files.make_dir(session_dir)
        except Exception:
            pass

        # Create a document
        doc_content = {
            "title": "Accessibility Test",
            "sections": [{"heading": "Test", "level": 1, "content": "Content here."}]
        }

        output_path = f"{session_dir}/accessible.docx"
        json_file = f"{session_dir}/content.json"
        sandbox.files.write(json_file, json.dumps(doc_content))

        result = sandbox.commands.run(
            f'python {SANDBOX_ROOT}/scripts/create_docx.py "{output_path}" --file "{json_file}"',
            timeout=60,
            cwd=session_dir,
        )

        assert result.exit_code == 0, f"Failed to create test document: {result.stderr}"

        yield {
            "manager": manager,
            "sandbox": sandbox,
            "user_id": test_user_id,
            "session_id": session_id,
            "file_path": output_path,
            "logical_path": "accessible.docx",
        }

        manager.close_sandbox(test_user_id)

    def test_file_exists_in_sandbox(self, created_docx):
        """Created DOCX file exists in sandbox filesystem."""
        sandbox = created_docx["sandbox"]
        file_path = created_docx["file_path"]

        # List directory and check file exists
        session_dir = str(Path(file_path).parent)
        entries = sandbox.files.list(session_dir)
        file_names = [e.name for e in entries]

        assert "accessible.docx" in file_names, f"File not found. Files: {file_names}"

    def test_file_can_be_read_as_bytes(self, created_docx):
        """DOCX file content can be read from sandbox using format='bytes'."""
        sandbox = created_docx["sandbox"]
        file_path = created_docx["file_path"]

        # CRITICAL: Use format="bytes" for binary files like DOCX
        content = sandbox.files.read(file_path, format="bytes")

        assert content is not None, "File content should not be None"
        assert len(content) > 0, "File should have content"

    def test_file_has_minimum_size(self, created_docx):
        """DOCX file has minimum expected size (not empty/corrupt)."""
        sandbox = created_docx["sandbox"]
        file_path = created_docx["file_path"]

        # CRITICAL: Use format="bytes" for binary files
        content = sandbox.files.read(file_path, format="bytes")

        size = len(content)

        # A valid DOCX should be at least 2KB (contains XML structure)
        assert size >= 2000, f"File too small ({size} bytes), likely corrupt"

    def test_bytes_format_returns_bytearray(self, created_docx):
        """format='bytes' returns bytearray for binary files."""
        sandbox = created_docx["sandbox"]
        file_path = created_docx["file_path"]

        # With format="bytes", E2B returns bytearray
        content = sandbox.files.read(file_path, format="bytes")

        assert isinstance(content, (bytes, bytearray)), f"Expected bytes/bytearray, got {type(content)}"

    def test_default_text_format_corrupts_binary(self, created_docx):
        """Default text format corrupts binary DOCX files - documents the bug."""
        sandbox = created_docx["sandbox"]
        file_path = created_docx["file_path"]

        # Default format="text" will corrupt binary data by decoding as UTF-8
        text_content = sandbox.files.read(file_path)  # default is text

        # This proves the issue: text mode returns string
        assert isinstance(text_content, str), "Default format returns string"

        # The string contains replacement characters because binary was decoded as UTF-8
        # This is why DOCX files appear corrupt when using default read()


class TestDocxRenderability:
    """Test 4: Verify DOCX files are valid and can be parsed."""

    @pytest.fixture
    def created_docx_content(self):
        """Create a DOCX and return its content."""
        manager = SandboxManager()
        test_user_id = f"test-user-{uuid.uuid4()}"
        session_id = str(uuid.uuid4())

        sandbox = manager.create_sandbox(test_user_id)
        session_dir = session_workspace_dir(session_id)

        try:
            sandbox.files.make_dir(session_dir)
        except Exception:
            pass

        doc_content = {
            "title": "Renderability Test",
            "sections": [
                {"heading": "Section 1", "level": 1, "content": "First section content."},
                {"heading": "Section 2", "level": 2, "bullet_points": ["A", "B", "C"]}
            ]
        }

        output_path = f"{session_dir}/renderable.docx"
        json_file = f"{session_dir}/content.json"
        sandbox.files.write(json_file, json.dumps(doc_content))

        result = sandbox.commands.run(
            f'python {SANDBOX_ROOT}/scripts/create_docx.py "{output_path}" --file "{json_file}"',
            timeout=60,
            cwd=session_dir,
        )

        assert result.exit_code == 0, f"Failed to create test document: {result.stderr}"

        # Read the content - CRITICAL: use format="bytes" for binary files
        content = sandbox.files.read(output_path, format="bytes")

        yield {
            "content": content,
            "expected_title": "Renderability Test",
            "expected_sections": ["Section 1", "Section 2"],
        }

        manager.close_sandbox(test_user_id)

    def test_valid_zip_structure(self, created_docx_content):
        """DOCX has valid ZIP structure (DOCX is a ZIP file)."""
        import zipfile

        content = created_docx_content["content"]

        # Content should already be bytes from format="bytes"
        buffer = BytesIO(bytes(content))

        try:
            with zipfile.ZipFile(buffer, 'r') as zf:
                # DOCX must contain these files
                names = zf.namelist()
                assert "[Content_Types].xml" in names, "Missing [Content_Types].xml"
                assert any("document.xml" in n for n in names), "Missing document.xml"
        except zipfile.BadZipFile as e:
            pytest.fail(f"Invalid ZIP structure: {e}")

    def test_parseable_by_python_docx(self, created_docx_content):
        """DOCX can be parsed by python-docx library."""
        from docx import Document

        content = created_docx_content["content"]

        # Content is already bytes from format="bytes"
        buffer = BytesIO(bytes(content))

        try:
            doc = Document(buffer)
            assert doc is not None, "Document should be parseable"
        except Exception as e:
            pytest.fail(f"python-docx failed to parse: {e}")

    def test_contains_expected_content(self, created_docx_content):
        """DOCX contains the expected text content."""
        from docx import Document

        content = created_docx_content["content"]
        expected_title = created_docx_content["expected_title"]

        buffer = BytesIO(bytes(content))
        doc = Document(buffer)

        # Extract all text from document
        full_text = "\n".join([p.text for p in doc.paragraphs])

        assert expected_title in full_text, f"Title '{expected_title}' not found in document"
        assert "Section 1" in full_text, "Section 1 heading not found"
        assert "First section content" in full_text, "Section 1 content not found"


class TestPptxGeneration:
    """PPTX pipeline: ensure presentations can be generated and parsed."""

    @pytest.fixture
    def created_pptx_content(self):
        """Create a PPTX file in the sandbox and return its content."""
        manager = SandboxManager()
        test_user_id = f"test-user-{uuid.uuid4()}"
        session_id = str(uuid.uuid4())

        sandbox = manager.create_sandbox(test_user_id)
        session_dir = session_workspace_dir(session_id)

        try:
            sandbox.files.make_dir(session_dir)
        except Exception:
            pass

        script_path = _ensure_html2pptx_script(sandbox)

        dep_check = sandbox.commands.run(
            'node -e "require(\'pptxgenjs\'); require(\'playwright\'); require(\'sharp\');"',
            timeout=60,
            cwd=SANDBOX_ROOT,
        )
        if dep_check.exit_code != 0:
            manager.close_sandbox(test_user_id)
            pytest.skip(f"Missing Node dependencies for html2pptx: {dep_check.stderr or dep_check.stdout}")

        slide1 = """<!doctype html>
<html>
<head>
  <style>
    body {
      width: 720pt;
      height: 405pt;
      margin: 0;
      padding: 32pt;
      display: flex;
      flex-direction: column;
      justify-content: center;
      gap: 16pt;
      font-family: Arial, sans-serif;
      background: #f8fafc;
    }
    h1 { margin: 0; font-size: 32pt; }
    p { margin: 0; font-size: 18pt; }
  </style>
</head>
<body>
  <h1>Welcome</h1>
  <p>This deck validates PPTX creation via html2pptx.</p>
</body>
</html>
"""
        slide2 = """<!doctype html>
<html>
<head>
  <style>
    body {
      width: 720pt;
      height: 405pt;
      margin: 0;
      padding: 32pt;
      display: flex;
      flex-direction: column;
      justify-content: flex-start;
      gap: 12pt;
      font-family: Arial, sans-serif;
      background: #eef2ff;
    }
    h2 { margin: 0; font-size: 28pt; }
    ul { margin: 0; padding-left: 28pt; font-size: 18pt; }
  </style>
</head>
<body>
  <h2>Agenda</h2>
  <ul>
    <li>Create file</li>
    <li>Read bytes</li>
    <li>Parse slides</li>
  </ul>
</body>
</html>
"""

        sandbox.files.write(f"{session_dir}/slide1.html", slide1)
        sandbox.files.write(f"{session_dir}/slide2.html", slide2)
        output_path = f"{session_dir}/pipeline.pptx"

        node_script = f"""
const path = require('path');
const pptxgen = require('pptxgenjs');
const html2pptx = require('{script_path}');

(async () => {{
  const pres = new pptxgen();
  pres.layout = 'LAYOUT_16x9';
  await html2pptx('slide1.html', pres, {{ tmpDir: '{SANDBOX_ROOT}/tmp' }});
  await html2pptx('slide2.html', pres, {{ tmpDir: '{SANDBOX_ROOT}/tmp' }});
  const outputPath = path.join(process.cwd(), 'pipeline.pptx');
  await pres.writeFile({{ fileName: outputPath }});
}})().catch(err => {{
  console.error(err);
  process.exit(1);
}});
"""

        sandbox.files.write(f"{session_dir}/generate_pptx.js", node_script)

        result = sandbox.commands.run(
            "node generate_pptx.js",
            timeout=120,
            cwd=session_dir,
        )

        assert result.exit_code == 0, f"Failed to create PPTX via html2pptx: {result.stderr or result.stdout}"

        content = sandbox.files.read(output_path, format="bytes")
        size = len(content)
        assert size > 2000, f"PPTX file too small ({size} bytes), likely corrupt"

        yield {
            "content": content,
            "expected_strings": ["Welcome", "Agenda", "Create file", "Read bytes", "Parse slides"],
        }

        manager.close_sandbox(test_user_id)

    def test_pptx_zip_structure(self, created_pptx_content):
        """PPTX should be a valid ZIP archive with slide contents."""
        import zipfile

        content = created_pptx_content["content"]
        buffer = BytesIO(bytes(content))

        try:
            with zipfile.ZipFile(buffer, 'r') as zf:
                names = zf.namelist()
                assert "[Content_Types].xml" in names, "Missing [Content_Types].xml"
                assert any(n.startswith("ppt/slides/slide") for n in names), "Missing slide XML files"
        except zipfile.BadZipFile as e:
            pytest.fail(f"Invalid PPTX ZIP structure: {e}")

    def test_pptx_parseable_by_python_pptx(self, created_pptx_content):
        """PPTX can be parsed by python-pptx."""
        from pptx import Presentation

        content = created_pptx_content["content"]
        prs = Presentation(BytesIO(bytes(content)))

        assert len(prs.slides) >= 2, "Expected at least two slides"

    def test_pptx_contains_expected_text(self, created_pptx_content):
        """PPTX slides contain expected text content."""
        from pptx import Presentation

        content = created_pptx_content["content"]
        expected_strings = created_pptx_content["expected_strings"]

        prs = Presentation(BytesIO(bytes(content)))

        slide_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    slide_text.append(shape.text)

        full_text = "\n".join(slide_text)
        for expected in expected_strings:
            assert expected in full_text, f"Slide text '{expected}' not found"


class TestEndToEndPipeline:
    """Integration test: Full pipeline from sandbox creation to valid document."""

    def test_full_pipeline(self):
        """Complete end-to-end test of document generation."""
        from docx import Document

        manager = SandboxManager()
        test_user_id = f"test-user-{uuid.uuid4()}"
        session_id = str(uuid.uuid4())

        try:
            # Step 1: Create sandbox
            print("\n[1/5] Creating sandbox...")
            sandbox = manager.create_sandbox(test_user_id)
            assert sandbox.is_running(), "Sandbox should be running"
            print(f"  ✓ Sandbox created: {sandbox.sandbox_id}")

            # Step 2: Create session workspace
            print("[2/5] Creating session workspace...")
            session_dir = session_workspace_dir(session_id)
            sandbox.files.make_dir(session_dir)
            print(f"  ✓ Session dir: {session_dir}")

            # Step 3: Create document
            print("[3/5] Creating DOCX document...")
            doc_content = {
                "title": "End-to-End Test Document",
                "sections": [
                    {
                        "heading": "Executive Summary",
                        "level": 1,
                        "content": "This document validates the entire pipeline."
                    },
                    {
                        "heading": "Key Points",
                        "level": 2,
                        "bullet_points": [
                            "Sandbox creation works",
                            "Document creation works",
                            "File reading works",
                            "Document parsing works"
                        ]
                    }
                ]
            }

            output_path = f"{session_dir}/e2e_test.docx"
            json_file = f"{session_dir}/e2e_content.json"
            sandbox.files.write(json_file, json.dumps(doc_content))

            result = sandbox.commands.run(
                f'python {SANDBOX_ROOT}/scripts/create_docx.py "{output_path}" --file "{json_file}"',
                timeout=60,
                cwd=session_dir,
            )
            assert result.exit_code == 0, f"Script failed: {result.stderr}"
            print(f"  ✓ Document created: {output_path}")

            # Step 4: Read and validate file
            print("[4/5] Reading document from sandbox...")
            # CRITICAL: Use format="bytes" for binary files like DOCX
            content = sandbox.files.read(output_path, format="bytes")
            assert content is not None
            assert len(content) > 2000, f"File too small: {len(content)} bytes"
            print(f"  ✓ Read {len(content)} bytes (as binary)")

            # Step 5: Parse document
            print("[5/5] Parsing document with python-docx...")
            # Content is already bytes/bytearray from format="bytes"
            buffer = BytesIO(bytes(content))
            doc = Document(buffer)

            # Verify content
            full_text = "\n".join([p.text for p in doc.paragraphs])
            assert "End-to-End Test Document" in full_text
            assert "Executive Summary" in full_text
            assert "Sandbox creation works" in full_text
            print("  ✓ Document parsed successfully")

            print("\n✅ All pipeline tests passed!")

        finally:
            manager.close_sandbox(test_user_id)


if __name__ == "__main__":
    pytest.main([__file__, "-v", "-s"])
